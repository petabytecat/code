import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
import re

# Increase the maximum image size limit
Image.MAX_IMAGE_PIXELS = None  # Remove image size restriction

def create_ranked_presentation(directory, output_file):
    # Get all PNG files in the directory
    files = [f for f in os.listdir(directory) if f.endswith('.png')]

    # Sort files by rank
    files.sort(key=lambda x: int(re.search(r'rank(\d+)', x).group(1)))

    # Create a new presentation
    prs = Presentation()

    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    for file in files:
        # Add a new slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Use the blank layout

        # Add the image to the slide
        img_path = os.path.join(directory, file)
        left = top = Inches(0)

        try:
            # Open image and get its size
            with Image.open(img_path) as img:
                width, height = img.size

            # Calculate scaling factor to fit slide
            scale_factor = min(prs.slide_width / width, prs.slide_height / height)

            # Add picture to slide
            slide.shapes.add_picture(img_path, left, top, width=width*scale_factor, height=height*scale_factor)
        except Exception as e:
            print(f"Error processing image {file}: {str(e)}")
            continue

        # Add filename as text box
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(14), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = file
        title_frame.paragraphs[0].font.size = Pt(14)

        # Extract rank from filename
        rank = re.search(r'rank(\d+)', file).group(1)

        # Add rank to top right corner
        rank_box = slide.shapes.add_textbox(Inches(14), Inches(0.5), Inches(1.5), Inches(0.5))
        rank_frame = rank_box.text_frame
        rank_frame.text = f"Rank: {rank}"
        rank_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        rank_frame.paragraphs[0].font.size = Pt(24)
        rank_frame.paragraphs[0].font.bold = True

    # Save the presentation
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

# Usage
directory = 'spec_sheets_pngs'  # Replace with the actual path to your folder
output_file = 'Ranked_Presentation.pptx'  # You can change the output filename if desired
create_ranked_presentation(directory, output_file)
